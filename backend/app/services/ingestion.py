import time
from pathlib import Path

from llama_index.core import Document
from llama_index.core.node_parser import SentenceSplitter
from llama_index.embeddings.huggingface import HuggingFaceEmbedding

from backend.app.core.config import PARSED_DIR
from backend.app.core.logging import setup_logger
from backend.app.services.vector_store import get_vector_store, store_embeddings

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

logger = setup_logger()


def parse_file(file_path: Path) -> str:
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        reader = PdfReader(str(file_path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if ext == ".docx":
        doc = DocxDocument(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)

    if ext == ".pptx":
        prs = Presentation(str(file_path))
        return "\n".join(
            shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
        )

    raise ValueError("Unsupported file format")


def ingest_document(
    file_path: Path,
    original_filename: str,
    document_id: str,
    session_id: str,  # 🔥 NEW
) -> None:
    logger.info(f"Parsing started | document_id={document_id}")

    start_time = time.time()
    parsed_text = parse_file(file_path)

    # Persist parsed text
    parsed_output_path = PARSED_DIR / f"{document_id}.txt"
    parsed_output_path.write_text(parsed_text, encoding="utf-8")

    # 1️⃣ Create Document
    document = Document(
        text=parsed_text,
        metadata={
            "document_id": document_id,
            "session_id": session_id,   # 🔥 CRITICAL
            "filename": original_filename,
        },
    )

    # 2️⃣ Chunk
    splitter = SentenceSplitter(chunk_size=512, chunk_overlap=50)
    nodes = splitter.get_nodes_from_documents([document])

    # 3️⃣ Embed
    embed_model = HuggingFaceEmbedding(
        model_name="nomic-ai/nomic-embed-text-v1.5",
        trust_remote_code=True,
    )
    for node in nodes:
        node.embedding = embed_model.get_text_embedding(node.text)

    # 4️⃣ Store in Milvus
    vector_store = get_vector_store()
    store_embeddings(nodes, vector_store)

    duration = round((time.time() - start_time) * 1000, 2)
    logger.info(f"Ingestion completed | document_id={document_id} | {duration}ms")
